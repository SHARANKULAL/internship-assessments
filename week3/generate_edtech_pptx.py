
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as ns
from lxml import etree
import copy

# ─── Color Palette ──────────────────────────────────────────────────────────────
BG_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE    = RGBColor(0x1A, 0x50, 0x9C)   # deep navy
ACCENT_TEAL    = RGBColor(0x00, 0x87, 0x8A)   # teal
ACCENT_AMBER   = RGBColor(0xF5, 0xA6, 0x23)   # amber
LIGHT_TEAL_BG  = RGBColor(0xE8, 0xF6, 0xF7)   # very light teal
LIGHT_BLUE_BG  = RGBColor(0xE9, 0xF1, 0xFB)   # very light blue
TEXT_DARK      = RGBColor(0x1A, 0x1A, 0x2E)   # near black
TEXT_GRAY      = RGBColor(0x55, 0x65, 0x7A)   # medium gray
DIVIDER_COLOR  = RGBColor(0x00, 0x87, 0x8A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ─── Helpers ────────────────────────────────────────────────────────────────────

def solid_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb

def add_rect(slide, l, t, w, h, rgb, radius=False):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # 1 = rectangle
        Inches(l), Inches(t), Inches(w), Inches(h))
    solid_fill(shape, rgb)
    shape.line.fill.background()
    if radius:
        # add rounded corners via XML
        sp = shape._element
        prstGeom = sp.find('.//' + ns.qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.attrib['prst'] = 'roundRect'
            avLst = prstGeom.find(ns.qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, ns.qn('a:avLst'))
            gd = etree.SubElement(avLst, ns.qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', 'val 30000')
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_name='Calibri', font_size=12, bold=False, italic=False,
                 color=TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, bullets, l, t, w, h,
                   font_name='Calibri', font_size=13, color=TEXT_DARK,
                   leading_char='▸  ', indent=0.15):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = leading_char + bullet
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        p.space_before = Pt(4)
        p.space_after  = Pt(4)
    return txBox

def add_pill(slide, text, l, t, w, h, bg_rgb, text_rgb=BG_WHITE, font_size=10):
    box = add_rect(slide, l, t, w, h, bg_rgb, radius=True)
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = 'Calibri'
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_rgb
    return box

def add_card(slide, title, body_lines, l, t, w, h, accent_rgb=ACCENT_BLUE):
    """Draws a card with top accent bar, title, and bullet body."""
    add_rect(slide, l, t, w, h, RGBColor(0xF4, 0xF7, 0xFD))      # card bg
    add_rect(slide, l, t, w, 0.06, accent_rgb)                     # top bar
    add_text_box(slide, title,
                 l+0.15, t+0.12, w-0.3, 0.35,
                 font_size=11, bold=True, color=accent_rgb)
    add_bullet_box(slide, body_lines,
                   l+0.15, t+0.50, w-0.3, h-0.65,
                   font_size=10, color=TEXT_GRAY, leading_char='• ')

# ─── Slide background helper ────────────────────────────────────────────────────

def set_slide_bg(slide):
    """White slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_WHITE

# ─── Left sidebar accent ────────────────────────────────────────────────────────

def add_left_accent(slide, rgb=ACCENT_BLUE):
    add_rect(slide, 0, 0, 0.25, 7.5, rgb)

def add_top_bar(slide, rgb=ACCENT_BLUE, height=0.08):
    add_rect(slide, 0, 0, 13.33, height, rgb)

def add_bottom_bar(slide, rgb=ACCENT_TEAL, height=0.06):
    add_rect(slide, 0, 7.5-height, 13.33, height, rgb)

def add_slide_number(slide, num, total=6):
    add_text_box(slide, f'{num} / {total}',
                 12.7, 7.2, 0.55, 0.25,
                 font_size=9, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)

def add_watermark_label(slide, label='EdTech DQA'):
    add_text_box(slide, label,
                 0.35, 7.22, 2.5, 0.22,
                 font_size=9, color=TEXT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title & Objective
# ═══════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(BLANK)
set_slide_bg(slide1)

# Left navy panel
add_rect(slide1, 0, 0, 4.5, 7.5, ACCENT_BLUE)
# Teal accent strip on panel
add_rect(slide1, 4.25, 0, 0.25, 7.5, ACCENT_TEAL)

# Badge / category pill
add_pill(slide1, 'STRATEGY PROPOSAL', 4.7, 0.5, 2.8, 0.38, ACCENT_TEAL, BG_WHITE, 9)

# Main title (right side)
add_text_box(slide1,
             'Data Quality Assessment\nStrategy for\nEdTech Platforms',
             4.7, 1.1, 8.2, 2.6,
             font_name='Calibri', font_size=36, bold=True,
             color=TEXT_DARK, align=PP_ALIGN.LEFT)

# Thin divider
add_rect(slide1, 4.7, 3.75, 5.0, 0.04, ACCENT_TEAL)

# Sub-heading
add_text_box(slide1,
             'Transforming Learning Outcomes Through\nData-Driven Excellence',
             4.7, 3.9, 8.1, 0.9,
             font_size=16, italic=True, color=ACCENT_TEAL, align=PP_ALIGN.LEFT)

# Purpose bullets (3-column icon-style metric cards)
metrics = [
    ('🎯', 'Assess', 'Evaluate current\ndata quality posture'),
    ('📊', 'Govern', 'Establish a structured\ndata governance model'),
    ('🚀', 'Grow', 'Enable data-driven\nstrategic decisions'),
]
for idx, (icon, title, body) in enumerate(metrics):
    cx = 4.7 + idx * 2.75
    add_rect(slide1, cx, 4.95, 2.55, 1.85, LIGHT_BLUE_BG, radius=True)
    add_text_box(slide1, icon,   cx+0.15, 5.05, 0.55, 0.55, font_size=22, color=ACCENT_BLUE)
    add_text_box(slide1, title,  cx+0.15, 5.55, 2.2,  0.3, font_size=11, bold=True, color=ACCENT_BLUE)
    add_text_box(slide1, body,   cx+0.15, 5.85, 2.2,  0.85, font_size=9.5, color=TEXT_GRAY)

# Key takeaway strip
add_rect(slide1, 4.7, 7.0, 8.55, 0.42, LIGHT_TEAL_BG)
add_text_box(slide1,
             '💡  Takeaway: Data quality directly impacts learning insights, operational efficiency, and business decisions.',
             4.85, 7.03, 8.3, 0.38,
             font_size=9.5, italic=True, color=ACCENT_TEAL)

# Left panel content
add_text_box(slide1, 'Prepared by', 0.3, 1.1, 3.8, 0.3, font_size=9, italic=True, color=RGBColor(0xCC, 0xDD, 0xFF))
add_text_box(slide1, 'Business Strategy\nConsulting Division', 0.3, 1.45, 3.8, 0.6, font_size=13, bold=True, color=BG_WHITE)
add_rect(slide1, 0.3, 2.15, 2.5, 0.04, ACCENT_AMBER)
add_text_box(slide1, 'March 2026', 0.3, 2.25, 3.8, 0.3, font_size=11, color=ACCENT_AMBER)
add_text_box(slide1, 'CONFIDENTIAL', 0.3, 7.12, 3.8, 0.28, font_size=8, bold=True, color=RGBColor(0xAA, 0xBB, 0xDD))
add_slide_number(slide1, 1)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Business Problem & Market Context
# ═══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(BLANK)
set_slide_bg(slide2)
add_top_bar(slide2, ACCENT_BLUE, 0.08)
add_bottom_bar(slide2, ACCENT_TEAL, 0.06)
add_left_accent(slide2, ACCENT_BLUE)

# Section label
add_text_box(slide2, 'SLIDE 02  |  BUSINESS PROBLEM & MARKET CONTEXT',
             0.4, 0.15, 12.0, 0.28, font_size=8.5, bold=True, color=ACCENT_TEAL)

# Slide title
add_text_box(slide2, 'Why Data Quality Matters Now',
             0.4, 0.55, 10.0, 0.65, font_size=28, bold=True, color=ACCENT_BLUE)

# Sub-title
add_text_box(slide2,
             'EdTech platforms are growing fast — but data fragmentation is holding them back.',
             0.4, 1.2, 12.5, 0.4, font_size=13, italic=True, color=TEXT_GRAY)

add_rect(slide2, 0.4, 1.66, 5.5, 0.04, ACCENT_TEAL)

# --- Left column: Problem statement cards ---
problems = [
    ('📈  Explosive Market Growth',
     ['Global EdTech market projected at $400B+ by 2028',
      'Rapid data volume growth outpaces governance maturity',
      'Quality gaps compound at scale']),
    ('🗂️  Data Fragmentation',
     ['Siloed systems: LMS, assessments, CRM, billing',
      'No single source of truth for learner journeys',
      'Inconsistent schema and field definitions across tools']),
    ('⚠️  Impact of Poor Quality',
     ['Misleading analytics misguide curriculum decisions',
      'Duplicate learner records distort engagement metrics',
      'Delayed data updates create stale dashboards']),
]
for i, (title, bullets) in enumerate(problems):
    ty = 1.8 + i * 1.7
    add_rect(slide2, 0.4, ty, 6.5, 1.55, LIGHT_BLUE_BG, radius=True)
    add_rect(slide2, 0.4, ty, 0.08, 1.55, ACCENT_BLUE)
    add_text_box(slide2, title, 0.6, ty+0.08, 6.1, 0.38, font_size=11.5, bold=True, color=ACCENT_BLUE)
    add_bullet_box(slide2, bullets, 0.6, ty+0.48, 6.1, 1.0, font_size=10, color=TEXT_GRAY, leading_char='→  ')

# --- Right column: stat highlights ---
stats = [
    ('73%', 'of EdTech companies\ncite data silos as a\ntop analytics barrier'),
    ('60%', 'reporting errors traced\nto duplicate or stale\nlearner records'),
    ('3x', 'faster decisions\nin data-mature\nEdTech firms'),
]
add_text_box(slide2, 'Industry Signals', 7.3, 1.8, 5.6, 0.4, font_size=13, bold=True, color=ACCENT_BLUE)
add_rect(slide2, 7.3, 2.22, 5.6, 0.04, ACCENT_AMBER)
for i, (num, label) in enumerate(stats):
    ty = 2.38 + i * 1.6
    add_rect(slide2, 7.3, ty, 5.6, 1.45, LIGHT_TEAL_BG, radius=True)
    add_text_box(slide2, num, 7.5, ty+0.1, 1.8, 0.9, font_size=36, bold=True, color=ACCENT_TEAL)
    add_text_box(slide2, label, 9.3, ty+0.22, 3.4, 0.9, font_size=10.5, color=TEXT_DARK)

# Takeaway
add_rect(slide2, 0.4, 7.0, 12.55, 0.4, LIGHT_TEAL_BG)
add_text_box(slide2,
             '💡  Takeaway: Poor data quality leads to inaccurate learning insights and operational inefficiencies.',
             0.6, 7.03, 12.2, 0.35, font_size=9.5, italic=True, color=ACCENT_TEAL)
add_slide_number(slide2, 2)
add_watermark_label(slide2)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Value Proposition
# ═══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(BLANK)
set_slide_bg(slide3)
add_top_bar(slide3, ACCENT_BLUE, 0.08)
add_bottom_bar(slide3, ACCENT_TEAL, 0.06)
add_left_accent(slide3, ACCENT_BLUE)

add_text_box(slide3, 'SLIDE 03  |  VALUE PROPOSITION',
             0.4, 0.15, 12.0, 0.28, font_size=8.5, bold=True, color=ACCENT_TEAL)
add_text_box(slide3, 'What a Data Quality Framework Delivers',
             0.4, 0.55, 12.0, 0.65, font_size=28, bold=True, color=ACCENT_BLUE)
add_text_box(slide3, 'From reactive firefighting to proactive data confidence.',
             0.4, 1.2, 12.0, 0.4, font_size=13, italic=True, color=TEXT_GRAY)
add_rect(slide3, 0.4, 1.66, 4.0, 0.04, ACCENT_TEAL)

# 4 value cards in 2x2 grid
value_cards = [
    ('📊', 'Reliable\nLearning Analytics',
     ['Accurate, consistent data feeds\ninto dashboards and reports',
      'Curriculum teams trust the numbers\nfor content decisions',
      'Reduces manual data reconciliation\nby 50–70%']),
    ('🎓', 'Better Student\nEngagement Insights',
     ['Unified learner profiles reveal\ntrue engagement patterns',
      'Early identification of at-risk students\nbefore dropout',
      'Personalised learning paths\nbased on clean cohort data']),
    ('⚙️', 'Improved Operational\nDecisions',
     ['Finance, HR, and operations rely on\nsingle-source validated data',
      'Faster monthly close cycles with\naccurate billing records',
      'Supports M&A due diligence and\ncompliance reporting']),
    ('✅', 'Reduced Reporting\nErrors & Rework',
     ['Automated validation cuts\nerror-correction labour',
      'Audit-ready data reduces regulatory\nrisk and fines',
      'Leadership confidence in board-level\nreporting increases']),
]
positions = [(0.4, 2.0), (6.7, 2.0), (0.4, 4.7), (6.7, 4.7)]
for (lft, top), (icon, title, bullet_list) in zip(positions, value_cards):
    add_rect(slide3, lft, top, 6.0, 2.55, LIGHT_BLUE_BG, radius=True)
    add_rect(slide3, lft, top, 6.0, 0.06, ACCENT_BLUE)
    add_text_box(slide3, icon,  lft+0.18, top+0.12, 0.7, 0.7, font_size=26, color=ACCENT_BLUE)
    add_text_box(slide3, title, lft+0.95, top+0.12, 4.8, 0.7, font_size=13, bold=True, color=ACCENT_BLUE)
    add_bullet_box(slide3, bullet_list,
                   lft+0.18, top+0.85, 5.65, 1.6,
                   font_size=10, color=TEXT_GRAY, leading_char='• ')

# Arrow connector hint (visual)
add_text_box(slide3, '⟶', 6.25, 3.15, 0.55, 0.55, font_size=20, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
add_text_box(slide3, '⟶', 6.25, 5.85, 0.55, 0.55, font_size=20, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

# Takeaway
add_rect(slide3, 0.4, 7.0, 12.55, 0.4, LIGHT_TEAL_BG)
add_text_box(slide3,
             '💡  Takeaway: A structured data quality framework enables better strategic decisions across the entire organisation.',
             0.6, 7.03, 12.2, 0.35, font_size=9.5, italic=True, color=ACCENT_TEAL)
add_slide_number(slide3, 3)
add_watermark_label(slide3)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Data Quality Framework (Business Perspective)
# ═══════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(BLANK)
set_slide_bg(slide4)
add_top_bar(slide4, ACCENT_BLUE, 0.08)
add_bottom_bar(slide4, ACCENT_TEAL, 0.06)
add_left_accent(slide4, ACCENT_BLUE)

add_text_box(slide4, 'SLIDE 04  |  DATA QUALITY FRAMEWORK',
             0.4, 0.15, 12.0, 0.28, font_size=8.5, bold=True, color=ACCENT_TEAL)
add_text_box(slide4, 'Proposed Data Quality Framework',
             0.4, 0.55, 12.0, 0.65, font_size=28, bold=True, color=ACCENT_BLUE)
add_text_box(slide4, 'A governance-led, business-aligned approach to data reliability.',
             0.4, 1.2, 12.0, 0.4, font_size=13, italic=True, color=TEXT_GRAY)
add_rect(slide4, 0.4, 1.66, 5.0, 0.04, ACCENT_TEAL)

# --- Framework Pillars (horizontal flow diagram) ---
pillars = [
    ('🏛️', 'Data\nGovernance',    'Define ownership,\npolicies & standards'),
    ('🔍', 'Profiling &\nAudit',   'Baseline quality\nassessment & scoring'),
    ('🤖', 'Automated\nValidation','Rule-based checks,\nalerts & auto-remediation'),
    ('🔁', 'Lifecycle\nManagement','Archival, retention\n& refresh policies'),
]
box_w, gap = 2.85, 0.22
start_x = 0.4
for i, (icon, title, desc) in enumerate(pillars):
    lx = start_x + i * (box_w + gap)
    add_rect(slide4, lx, 2.0, box_w, 2.2, ACCENT_BLUE, radius=True)
    add_text_box(slide4, icon, lx+0.15, 2.1, 0.9, 0.75, font_size=28, color=BG_WHITE)
    add_text_box(slide4, title, lx+0.15, 2.85, box_w-0.3, 0.6, font_size=12, bold=True, color=BG_WHITE)
    add_text_box(slide4, desc,  lx+0.15, 3.45, box_w-0.3, 0.65, font_size=10, color=RGBColor(0xCC, 0xDD, 0xFF))

# Arrows between pillars
for i in range(3):
    ax = start_x + (i+1)*(box_w+gap) - gap - 0.05
    add_text_box(slide4, '▶', ax, 2.95, 0.28, 0.35, font_size=14, color=ACCENT_AMBER)

# --- Governance Roles Table ---
add_text_box(slide4, 'Governance Accountability Model',
             0.4, 4.42, 8.0, 0.38, font_size=13, bold=True, color=ACCENT_BLUE)
roles = [
    ('Chief Data Officer',       'Strategy & Policy oversight',       'Enterprise-wide'),
    ('Data Stewards',            'Domain-level quality ownership',     'Per product/LMS module'),
    ('Data Engineering Team',    'Automated validation pipelines',     'Technical layer'),
    ('Business Analysts',        'KPI reporting & insight validation', 'Analytics & BI tools'),
]
headers = ('Role', 'Responsibility', 'Scope')
col_w = [3.2, 4.8, 4.5]
col_x = [0.4, 3.65, 8.5]
row_h = 0.34
# header row
for ci, hdr in enumerate(headers):
    add_rect(slide4, col_x[ci], 4.85, col_w[ci]-0.08, row_h, ACCENT_BLUE)
    add_text_box(slide4, hdr, col_x[ci]+0.1, 4.87, col_w[ci]-0.18, 0.3,
                 font_size=10, bold=True, color=BG_WHITE)
for ri, row in enumerate(roles):
    ty = 4.85 + (ri+1)*row_h
    bg = LIGHT_BLUE_BG if ri % 2 == 0 else BG_WHITE
    for ci, cell in enumerate(row):
        add_rect(slide4, col_x[ci], ty, col_w[ci]-0.08, row_h, bg)
        add_text_box(slide4, cell, col_x[ci]+0.1, ty+0.04, col_w[ci]-0.18, 0.28,
                     font_size=9.5, color=TEXT_DARK)

# Takeaway
add_rect(slide4, 0.4, 7.0, 12.55, 0.4, LIGHT_TEAL_BG)
add_text_box(slide4,
             '💡  Takeaway: Governance structures and automated checks together ensure long-term, scalable data reliability.',
             0.6, 7.03, 12.2, 0.35, font_size=9.5, italic=True, color=ACCENT_TEAL)
add_slide_number(slide4, 4)
add_watermark_label(slide4)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Business Model & Implementation Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(BLANK)
set_slide_bg(slide5)
add_top_bar(slide5, ACCENT_BLUE, 0.08)
add_bottom_bar(slide5, ACCENT_TEAL, 0.06)
add_left_accent(slide5, ACCENT_BLUE)

add_text_box(slide5, 'SLIDE 05  |  IMPLEMENTATION APPROACH',
             0.4, 0.15, 12.0, 0.28, font_size=8.5, bold=True, color=ACCENT_TEAL)
add_text_box(slide5, 'Phased Implementation Roadmap',
             0.4, 0.55, 12.0, 0.65, font_size=28, bold=True, color=ACCENT_BLUE)
add_text_box(slide5, 'Data quality as a strategic organisational capability — built iteratively.',
             0.4, 1.2, 12.0, 0.4, font_size=13, italic=True, color=TEXT_GRAY)
add_rect(slide5, 0.4, 1.66, 5.5, 0.04, ACCENT_TEAL)

# Phase roadmap (horizontal swimlane)
phases = [
    ('Phase 1', '0–4 Weeks',
     'Discovery &\nAssessment',
     ['Data inventory & profiling',
      'Stakeholder interviews',
      'Quality gap analysis',
      'Risk scoring baseline']),
    ('Phase 2', '5–10 Weeks',
     'Framework\nDesign',
     ['Governance charter',
      'Quality rules library',
      'KPI & metric definitions',
      'Tool & platform selection']),
    ('Phase 3', '11–18 Weeks',
     'Pilot\nImplementation',
     ['Automated checks deployed',
      'Monitoring dashboards live',
      'Training for data stewards',
      'Quick-win wins published']),
    ('Phase 4', '19–26 Weeks',
     'Scale &\nOptimise',
     ['Enterprise-wide rollout',
      'Advanced anomaly detection',
      'Continuous improvement cycle',
      'ROI measurement & reporting']),
]
phase_colors = [ACCENT_BLUE, ACCENT_TEAL, RGBColor(0x21, 0x8B, 0xC9), ACCENT_AMBER]
pw = 2.95
px_start = 0.4
for i, (phase, timeline, title, tasks) in enumerate(phases):
    px = px_start + i * (pw + 0.11)
    # phase header
    add_rect(slide5, px, 2.0, pw, 0.52, phase_colors[i], radius=True)
    add_text_box(slide5, phase,    px+0.12, 2.04, pw-0.24, 0.24, font_size=10, bold=True, color=BG_WHITE)
    add_text_box(slide5, timeline, px+0.12, 2.28, pw-0.24, 0.22, font_size=9,  color=RGBColor(0xCC, 0xEE, 0xFF))
    # title bar
    add_rect(slide5, px, 2.55, pw, 0.45, LIGHT_BLUE_BG)
    add_text_box(slide5, title, px+0.12, 2.57, pw-0.24, 0.42, font_size=11, bold=True, color=phase_colors[i])
    # tasks
    add_bullet_box(slide5, tasks, px+0.12, 3.06, pw-0.24, 1.9,
                   font_size=9.5, color=TEXT_DARK, leading_char='✓  ')
    # connector arrows
    if i < 3:
        ax = px + pw + 0.02
        add_text_box(slide5, '▶', ax, 2.62, 0.15, 0.35, font_size=11, color=ACCENT_AMBER)

# Integration note
add_rect(slide5, 0.4, 5.22, 12.55, 1.6, LIGHT_TEAL_BG, radius=True)
add_text_box(slide5, '🔗  Integration Touchpoints',
             0.6, 5.32, 5.0, 0.35, font_size=12, bold=True, color=ACCENT_TEAL)
integrations = [
    '📚  LMS & Course Management Systems (Moodle, Canvas, Blackboard)',
    '📊  BI & Analytics Platforms (Power BI, Tableau, Looker)',
    '💳  Payment & CRM Data Pipelines (Salesforce, Stripe)',
    '🛡️  Regulatory Compliance & Audit Systems',
]
add_bullet_box(slide5, integrations, 0.6, 5.7, 12.0, 1.0,
               font_size=10.5, color=TEXT_DARK, leading_char='')

# Takeaway
add_rect(slide5, 0.4, 7.0, 12.55, 0.4, LIGHT_TEAL_BG)
add_text_box(slide5,
             '💡  Takeaway: A phased approach minimises disruption while progressively building data-driven decision-making capabilities.',
             0.6, 7.03, 12.2, 0.35, font_size=9.5, italic=True, color=ACCENT_TEAL)
add_slide_number(slide5, 5)
add_watermark_label(slide5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Success Metrics & Business Impact
# ═══════════════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(BLANK)
set_slide_bg(slide6)
add_top_bar(slide6, ACCENT_BLUE, 0.08)
add_bottom_bar(slide6, ACCENT_TEAL, 0.06)
add_left_accent(slide6, ACCENT_BLUE)

add_text_box(slide6, 'SLIDE 06  |  SUCCESS METRICS & BUSINESS IMPACT',
             0.4, 0.15, 12.0, 0.28, font_size=8.5, bold=True, color=ACCENT_TEAL)
add_text_box(slide6, 'Measuring What Matters',
             0.4, 0.55, 12.0, 0.65, font_size=28, bold=True, color=ACCENT_BLUE)
add_text_box(slide6, 'Quantifiable outcomes that demonstrate return on the data quality investment.',
             0.4, 1.2, 12.0, 0.4, font_size=13, italic=True, color=TEXT_GRAY)
add_rect(slide6, 0.4, 1.66, 5.5, 0.04, ACCENT_TEAL)

# KPI Cards Row 1 (key numbers)
kpis = [
    ('40–60%', 'Reduction in\nData Errors',     ACCENT_BLUE),
    ('35%',    'Improvement in\nReporting Accuracy', ACCENT_TEAL),
    ('50%',    'Faster Analytics\nGeneration',  RGBColor(0x21, 0x8B, 0xC9)),
    ('25%',    'Operational\nEfficiency Gain',  ACCENT_AMBER),
]
kw = 2.95
kx_start = 0.4
for i, (num, label, col) in enumerate(kpis):
    kx = kx_start + i*(kw+0.12)
    add_rect(slide6, kx, 1.85, kw, 1.55, LIGHT_BLUE_BG, radius=True)
    add_rect(slide6, kx, 1.85, kw, 0.06, col)
    add_text_box(slide6, num,   kx+0.18, 1.95, kw-0.36, 0.7, font_size=30, bold=True, color=col)
    add_text_box(slide6, label, kx+0.18, 2.65, kw-0.36, 0.65, font_size=10.5, color=TEXT_DARK)

# Impact narrative cards (3 columns)
impact_cards = [
    ('📈  Reporting &\nAnalytics Accuracy',
     ['Board-level dashboards reflect real-time,\nvalidated data',
      'Curriculum insights are reliable enough to\ndrive investment decisions',
      'Audit and regulatory reports produced\nin hours, not days']),
    ('🏫  Student Success\n& Engagement',
     ['At-risk learner identification improves\nby up to 30%',
      'Personalised nudges based on clean\nbehavioural data increase retention',
      'Instructor workload reduced through\nautomated grade reconciliation']),
    ('💼  Strategic &\nOperational ROI',
     ['Data trust reduces over-verification\nmeetings and rework cycles',
      'Faster time-to-insight enables agile\nproduct and pricing decisions',
      'Investor-grade data readiness supports\nfundraising and M&A activity']),
]
ic_w = 4.0
ic_x_start = 0.4
for i, (title, bullets) in enumerate(impact_cards):
    ix = ic_x_start + i*(ic_w+0.16)
    add_rect(slide6, ix, 3.6, ic_w, 3.25, LIGHT_TEAL_BG, radius=True)
    add_rect(slide6, ix, 3.6, ic_w, 0.06, [ACCENT_BLUE, ACCENT_TEAL, ACCENT_AMBER][i])
    add_text_box(slide6, title, ix+0.18, 3.72, ic_w-0.36, 0.7,
                 font_size=12, bold=True, color=ACCENT_BLUE)
    add_bullet_box(slide6, bullets, ix+0.18, 4.45, ic_w-0.36, 2.3,
                   font_size=9.5, color=TEXT_DARK, leading_char='• ')

# CTA box
add_rect(slide6, 0.4, 7.0, 12.55, 0.4, LIGHT_TEAL_BG)
add_text_box(slide6,
             '💡  Takeaway: Data quality drives measurable improvements in business performance — the ROI is both quantifiable and strategic.',
             0.6, 7.03, 12.2, 0.35, font_size=9.5, italic=True, color=ACCENT_TEAL)
add_slide_number(slide6, 6)
add_watermark_label(slide6)

# ─── Save ────────────────────────────────────────────────────────────────────────
output = r'C:\Users\shara\OneDrive\Desktop\Aptitude\week3\EdTech_DataQuality_Deck.pptx'
prs.save(output)
print(f'✅  Presentation saved: {output}')
