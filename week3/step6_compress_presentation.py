"""
step6_compress_presentation.py
Takes the 8-slide EdTech_DataQuality_Deck.pptx and removes less critical slides
(Value Proposition, Implementation Roadmap, Success Metrics)
to fit a 5-minute presentation format.
"""

from pptx import Presentation

FILE_PATH = r'C:\Users\shara\OneDrive\Desktop\Aptitude\week3\EdTech_DataQuality_Deck.pptx'

prs = Presentation(FILE_PATH)

# We have 8 slides currently:
# 0 (Slide 1): Title & Objective
# 1 (Slide 2): Business Problem & Market Context
# 2 (Slide 3): Value Proposition  <-- Remove
# 3 (Slide 4): Data Quality Framework
# 4 (Slide 5): Implementation Roadmap  <-- Remove
# 5 (Slide 6): Success Metrics & Business Impact  <-- Remove
# 6 (Slide 7): Dataset Structure & Issues Found
# 7 (Slide 8): Data Cleaning Results

# To delete slides in python-pptx, we grab the slide's XML element and remove it.
# We must delete from highest index to lowest so indices don't shift.

slides_to_delete = [5, 4, 2] # corresponding to Slides 6, 5, 3

for index in slides_to_delete:
    rId = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]

# Re-number the remaining slides (Slide 1-5)
for i, slide in enumerate(prs.slides):
    # Find the slide number text box and update it
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if ' / 8' in run.text or ' / 6' in run.text:
                    run.text = f'{i+1} / 5'
                if 'SLIDE 0' in run.text:
                    # e.g. "SLIDE 04  |  DATA QUALITY FRAMEWORK" update to "SLIDE 03 ..."
                    # The parts are split, let's just do a simple replace on the full text
                    pass # We'll do a simpler update on the whole text frame

    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        # Update the breadcrumb label (e.g., SLIDE 04 | ...)
        if 'SLIDE 0' in shape.text:
            text = shape.text
            parts = text.split('  |  ')
            if len(parts) == 2:
                shape.text = f'SLIDE {i+1:02d}  |  {parts[1]}'

prs.save(FILE_PATH)
print("✅ Presentation compressed to 5 essential slides.")
