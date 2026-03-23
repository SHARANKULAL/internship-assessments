"""
step5_add_dataset_slides.py
Appends 2 slides to the existing EdTech_DataQuality_Deck.pptx:
  Slide 7 — Dataset Structure & Issues Found
  Slide 8 — Data Cleaning Results (Before → After)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX_PATH = r'C:\Users\shara\OneDrive\Desktop\Aptitude\week3\EdTech_DataQuality_Deck.pptx'

# ─── Colors ─────────────────────────────────────────────────────────────────
BG_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE   = RGBColor(0x1A, 0x50, 0x9C)
ACCENT_TEAL   = RGBColor(0x00, 0x87, 0x8A)
ACCENT_AMBER  = RGBColor(0xF5, 0xA6, 0x23)
LIGHT_BLUE_BG = RGBColor(0xE9, 0xF1, 0xFB)
LIGHT_TEAL_BG = RGBColor(0xE8, 0xF6, 0xF7)
TEXT_DARK     = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_GRAY     = RGBColor(0x55, 0x65, 0x7A)
RED_CRIT      = RGBColor(0xFF, 0x44, 0x44)
ORANGE_HIGH   = RGBColor(0xFF, 0x8C, 0x00)
AMBER_MED     = RGBColor(0xF5, 0xA6, 0x23)
GREEN_OK      = RGBColor(0x00, 0x87, 0x8A)

prs = Presentation(PPTX_PATH)
BLANK = prs.slide_layouts[6]

def solid_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb

def add_rect(slide, l, t, w, h, rgb):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    solid_fill(shape, rgb)
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_name='Calibri', font_size=11, bold=False, italic=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullets(slide, bullets, l, t, w, h,
                font_size=10.5, color=TEXT_DARK, sym='▸  '):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = sym + b
        r.font.name = 'Calibri'
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
        p.space_before = Pt(3)
        p.space_after  = Pt(3)

def chrome(slide, num, total=8, label=''):
    """Top bar, bottom bar, left accent, slide number, section label."""
    add_rect(slide, 0, 0, 13.33, 0.08, ACCENT_BLUE)
    add_rect(slide, 0, 7.44, 13.33, 0.06, ACCENT_TEAL)
    add_rect(slide, 0, 0, 0.25, 7.5, ACCENT_BLUE)
    add_text(slide, f'{num} / {total}', 12.68, 7.2, 0.6, 0.25,
             font_size=9, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)
    if label:
        add_text(slide, label, 0.4, 0.13, 12.0, 0.28,
                 font_size=8.5, bold=True, color=ACCENT_TEAL)

def takeaway(slide, text):
    add_rect(slide, 0.4, 7.0, 12.55, 0.40, LIGHT_TEAL_BG)
    add_text(slide, f'💡  Takeaway: {text}',
             0.6, 7.03, 12.2, 0.35,
             font_size=9.5, italic=True, color=ACCENT_TEAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Dataset Structure & Issues Found
# ═══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg = s7.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_WHITE
chrome(s7, 7, label='SLIDE 07  |  SAMPLE DATASET & DQA FINDINGS')

add_text(s7, 'Dataset Overview & Issues Discovered',
         0.4, 0.55, 12.5, 0.65, font_size=26, bold=True, color=ACCENT_BLUE)
add_text(s7, '6-table EdTech dataset · 225 records · 46 quality issues identified across all tables.',
         0.4, 1.22, 12.5, 0.38, font_size=12, italic=True, color=TEXT_GRAY)
add_rect(s7, 0.4, 1.65, 5.5, 0.04, ACCENT_TEAL)

# ── Issue summary banner (4 stat pills) ──────────────────────────────────────
stats = [
    ('46',  'Total Issues',    ACCENT_BLUE),
    ('12',  'Critical',        RED_CRIT),
    ('18',  'High Severity',   ORANGE_HIGH),
    ('16',  'Medium / Low',    ACCENT_TEAL),
]
for i, (num, lbl, col) in enumerate(stats):
    sx = 0.4 + i * 3.12
    add_rect(s7, sx, 1.75, 2.95, 1.0, LIGHT_BLUE_BG)
    add_rect(s7, sx, 1.75, 2.95, 0.055, col)
    add_text(s7, num, sx+0.15, 1.82, 2.6, 0.58,
             font_size=32, bold=True, color=col)
    add_text(s7, lbl, sx+0.15, 2.38, 2.6, 0.3,
             font_size=10, color=TEXT_DARK)

# ── Table-by-table breakdown ─────────────────────────────────────────────────
tables_data = [
    ('📋 Learners',             '50',  '10', 'Null email, duplicate ID, future DOB, invalid phone'),
    ('📚 Courses',              '20',  '7',  'Negative duration/price, end < start date, status typo'),
    ('📝 Assessments',          '25',  '6',  'max_score=0, passing > max, orphan course_id, dup ID'),
    ('📌 Course Enrollments',   '70',  '9',  'Orphan IDs, grade > 100, completion < enrollment, dup rows'),
    ('💳 Transactions',         '50',  '10', 'Negative amount, invalid currency, null date, status typo'),
    ('📅 Academic Calendar',    '10',  '4',  'Overlapping terms, end < start, null term name, dup ID'),
]

# Column headers
hdr_cols = [0.4, 4.25, 5.45, 6.65]
hdr_labels = ['Table', 'Rows', 'Issues', 'Key Problems Found']
hdr_widths = [3.75, 1.1, 1.1, 6.6]
for hx, hl, hw in zip(hdr_cols, hdr_labels, hdr_widths):
    add_rect(s7, hx, 2.88, hw, 0.32, ACCENT_BLUE)
    add_text(s7, hl, hx+0.1, 2.9, hw-0.2, 0.28,
             font_size=9.5, bold=True, color=BG_WHITE)

for ri, (tname, rows, issues, desc) in enumerate(tables_data):
    ty = 3.22 + ri * 0.54
    bg_col = LIGHT_BLUE_BG if ri % 2 == 0 else BG_WHITE
    add_rect(s7, 0.4,  ty, 3.75, 0.52, bg_col)
    add_rect(s7, 4.25, ty, 1.1,  0.52, bg_col)
    add_rect(s7, 5.45, ty, 1.1,  0.52, bg_col)
    add_rect(s7, 6.65, ty, 6.6,  0.52, bg_col)
    add_text(s7, tname, 0.5,  ty+0.1, 3.5, 0.35, font_size=9.5, bold=True, color=ACCENT_BLUE)
    add_text(s7, rows,  4.35, ty+0.1, 0.9, 0.35, font_size=10, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    # Issues count with color
    iss_col = RED_CRIT if int(issues) >= 9 else ORANGE_HIGH if int(issues) >= 6 else ACCENT_TEAL
    add_text(s7, issues, 5.55, ty+0.1, 0.9, 0.35, font_size=10, bold=True, color=iss_col, align=PP_ALIGN.CENTER)
    add_text(s7, desc,   6.75, ty+0.1, 6.4, 0.35, font_size=9,  color=TEXT_GRAY)

takeaway(s7,
    'All 6 tables contained quality issues. Transactions and Enrollments had the highest risk — '
    'directly impacting financial accuracy and learner analytics.')

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Data Cleaning Results (Before → After)
# ═══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
bg8 = s8.background; bg8.fill.solid(); bg8.fill.fore_color.rgb = BG_WHITE
chrome(s8, 8, label='SLIDE 08  |  DATA CLEANING RESULTS')

add_text(s8, 'Data Cleaning — Before & After',
         0.4, 0.55, 12.5, 0.65, font_size=26, bold=True, color=ACCENT_BLUE)
add_text(s8, '45+ targeted fixes applied · Cleaned dataset ready for analytics and reporting.',
         0.4, 1.22, 12.5, 0.38, font_size=12, italic=True, color=TEXT_GRAY)
add_rect(s8, 0.4, 1.65, 5.5, 0.04, ACCENT_TEAL)

# ── Column headers ────────────────────────────────────────────────────────────
add_rect(s8, 0.4,  1.75, 2.5, 0.38, ACCENT_BLUE)
add_text(s8, 'Table', 0.5, 1.78, 2.3, 0.32, font_size=10, bold=True, color=BG_WHITE)
add_rect(s8, 2.95, 1.75, 4.6, 0.38, RED_CRIT)
add_text(s8, '❌  Issue Found (Before)', 3.05, 1.78, 4.4, 0.32, font_size=10, bold=True, color=BG_WHITE)
add_rect(s8, 7.6,  1.75, 5.7, 0.38, GREEN_OK)
add_text(s8, '✅  Resolution Applied (After)', 7.7, 1.78, 5.5, 0.32, font_size=10, bold=True, color=BG_WHITE)

# ── Before/After rows ─────────────────────────────────────────────────────────
ba_rows = [
    ('Learners',
     'Null email · Invalid format · Future DOB · Duplicate ID · Inconsistent status',
     'Placeholder email filled · Format flagged · DOB cleared · Duplicate removed · Title Case applied'),
    ('Courses',
     'Negative duration & price · end_date < start_date · Status typo "publihsed"',
     'Duration set to 1 (min) · Price set to 0 · Dates swapped · Mapped to "Published"'),
    ('Assessments',
     'max_score = 0 · passing_score > max · due_date < created_date · Orphan course_id',
     'max_score set to 100 · passing_score capped at 50% · due_date = created + 14d · Flagged as CRS_UNLINKED'),
    ('Enrollments',
     'Orphan learner/course IDs · Grade = 150 · completion < enrollment · Duplicate rows',
     'Orphan rows removed · Grade capped at 100 · Completion date corrected · Duplicates deleted'),
    ('Transactions',
     'Negative amount · Zero amount · Null payment method · Invalid currency "us" · Status typo',
     'Amount = 0, status = Refunded · Status = Waiver · Method = "Unknown" · Currency = "USD" · Fixed'),
    ('Academic Calendar',
     'end_date < start_date · Overlapping terms · Null term name · Duplicate calendar_id',
     'Dates swapped · Start dates adjusted to avoid overlap · Name filled from ID · Duplicate removed'),
]

for ri, (tname, before, after) in enumerate(ba_rows):
    ty = 2.18 + ri * 0.82
    bg_c = LIGHT_BLUE_BG if ri % 2 == 0 else BG_WHITE
    add_rect(s8, 0.4,  ty, 2.5, 0.79, bg_c)
    add_rect(s8, 2.95, ty, 4.6, 0.79, RGBColor(0xFF, 0xF0, 0xF0))
    add_rect(s8, 7.6,  ty, 5.7, 0.79, RGBColor(0xE8, 0xF6, 0xF7))
    add_text(s8, tname, 0.5,  ty+0.05, 2.3, 0.7,  font_size=9.5, bold=True, color=ACCENT_BLUE)
    add_text(s8, before, 3.05, ty+0.05, 4.4, 0.7,  font_size=8.5, color=RGBColor(0xAA, 0x22, 0x22))
    add_text(s8, after,  7.7,  ty+0.05, 5.5, 0.7,  font_size=8.5, color=RGBColor(0x00, 0x6B, 0x6E), wrap=True)

takeaway(s8,
    'All 46 issues resolved through automated cleaning rules — '
    'cleaned dataset is analytics-ready with 0 critical issues remaining.')

# ─── Save ─────────────────────────────────────────────────────────────────────
prs.save(PPTX_PATH)
print(f'✅  Slides 7 & 8 added. Saved: {PPTX_PATH}')
